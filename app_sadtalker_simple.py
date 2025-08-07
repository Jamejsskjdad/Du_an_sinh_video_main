import os, sys
import gradio as gr
from src.gradio_demo import SadTalker  
from gtts import gTTS
from datetime import datetime
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from moviepy.editor import VideoFileClip, concatenate_videoclips, CompositeVideoClip, TextClip, ImageClip
from PIL import Image
import re
import json
from pptx import Presentation
import io
import time
import shutil

try:
    import webui  # in webui
    in_webui = True
except:
    in_webui = False

def toggle_audio_file(choice):
    if choice == False:
        return gr.update(visible=True), gr.update(visible=False)
    else:
        return gr.update(visible=False), gr.update(visible=True)
    
def ref_video_fn(path_of_ref_video):
    if path_of_ref_video is not None:
        return gr.update(value=True)
    else:
        return gr.update(value=False)

def read_text_file(file):
    """
    Read text content from uploaded file
    """
    if file is None:
        return "", "❌ Vui lòng chọn file văn bản!"
    
    try:
        # Get file extension
        file_path = file.name
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Check if it's a supported file type
        if file_ext not in ['.txt', '.md', '.doc', '.docx', '.rtf']:
            return "", "❌ Chỉ hỗ trợ file văn bản (.txt, .md, .doc, .docx, .rtf)"
        
        content = ""
        
        # Handle different file types
        if file_ext == '.docx':
            # For .docx files, we need to extract text properly
            try:
                import zipfile
                import xml.etree.ElementTree as ET
                
                # .docx is a ZIP file containing XML
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    # Find the document.xml file
                    if 'word/document.xml' in zip_file.namelist():
                        xml_content = zip_file.read('word/document.xml')
                        root = ET.fromstring(xml_content)
                        
                        # Extract text from all text elements
                        text_elements = []
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                text_elements.append(elem.text.strip())
                        
                        content = ' '.join(text_elements)
                    else:
                        return "", "❌ Không thể đọc nội dung từ file .docx"
                        
            except Exception as e:
                return "", f"❌ Lỗi đọc file .docx: {str(e)}"
        
        elif file_ext == '.doc':
            # For .doc files, we'll try to read as text but warn user
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        content = f.read()
                except:
                    return "", "❌ Không thể đọc file .doc. Vui lòng chuyển đổi sang .txt hoặc .docx"
        
        else:
            # For .txt, .md, .rtf files
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
        
        # Clean content - remove null bytes and other control characters
        content = ''.join(char for char in content if ord(char) >= 32 or char in '\n\r\t')
        
        # Strip whitespace and check if empty
        content = content.strip()
        if not content:
            return "", "❌ File văn bản trống hoặc chỉ chứa ký tự đặc biệt!"
        
        return content, f"✅ Đã đọc thành công file: {os.path.basename(file_path)} ({len(content)} ký tự thực tế)"
    
    except Exception as e:
        return "", f"❌ Lỗi đọc file: {str(e)}"

def convert_text_to_audio(text, language='vi'):
    """
    Convert text to audio using gTTS
    Returns the path to the generated audio file
    """
    try:
        if not text or text.strip() == "":
            return None
        
        # Create a temporary file for the audio
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
        temp_path = temp_file.name
        temp_file.close()
        
        # Convert text to speech
        tts = gTTS(text=text, lang=language, slow=False)
        tts.save(temp_path)
        
        # Return the audio file path directly for Gradio Audio component
        return temp_path
    except Exception as e:
        return None

def extract_text_from_pptx(pptx_file):
    """
    Extract text from each slide in PowerPoint file
    Returns list of text content from each slide
    """
    try:
        slides_text = []
        
        with zipfile.ZipFile(pptx_file.name, 'r') as zip_file:
            # Get all slide files
            slide_files = [f for f in zip_file.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            slide_files.sort()  # Sort by slide number
            
            for slide_file in slide_files:
                xml_content = zip_file.read(slide_file)
                root = ET.fromstring(xml_content)
                
                # Extract text from all text elements in the slide
                text_elements = []
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        text_elements.append(elem.text.strip())
                
                slide_text = ' '.join(text_elements)
                if slide_text.strip():
                    slides_text.append(slide_text)
        
        return slides_text
    except Exception as e:
        return []

def extract_slides_from_pptx(pptx_file):
    """
    Extract slides as images from PowerPoint file
    Returns list of slide images and their text content
    """
    try:
        slides_data = []
        prs = Presentation(pptx_file.name)
        
        for i, slide in enumerate(prs.slides):
            # Extract text from slide
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            
            # Convert slide to image (simplified approach)
            # In a real implementation, you'd use a library like python-pptx with rendering
            # For now, we'll create a placeholder image with text
            slide_text = slide_text.strip()
            if slide_text:
                slides_data.append({
                    'slide_number': i + 1,
                    'text': slide_text,
                    'image_path': None  # Will be created later
                })
        
        return slides_data
    except Exception as e:
        print(f"Error extracting slides: {str(e)}")
        return []

def get_audio_duration(audio_path):
    """
    Get duration of audio file in seconds
    """
    try:
        if audio_path and os.path.exists(audio_path):
            # Use AudioFileClip instead of VideoFileClip for audio files
            from moviepy.editor import AudioFileClip
            audio_clip = AudioFileClip(audio_path)
            duration = audio_clip.duration
            audio_clip.close()
            return duration
        return 0
    except Exception as e:
        print(f"Error getting audio duration: {str(e)}")
        return 0

def create_slide_image_with_text(text, output_path, width=1920, height=1080):
    """
    Create a slide image with text (placeholder for actual slide rendering)
    """
    try:
        # Create a white background image
        img = Image.new('RGB', (width, height), color='white')
        
        # For now, we'll create a simple text image
        # In a real implementation, you'd render the actual PowerPoint slide
        from PIL import ImageDraw, ImageFont
        
        draw = ImageDraw.Draw(img)
        
        # Try to use a default font, fallback to basic if not available
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except:
            font = ImageFont.load_default()
        
        # Draw text in the center
        text_lines = text.split('\n')
        y_position = height // 2 - (len(text_lines) * 50) // 2
        
        for line in text_lines:
            # Get text size
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            
            # Calculate position to center text
            x_position = (width - text_width) // 2
            
            # Draw text
            draw.text((x_position, y_position), line, fill='black', font=font)
            y_position += text_height + 20
        
        # Save image
        img.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error creating slide image: {str(e)}")
        return None

def create_lecture_video(sad_talker, slides_data, source_image, language, preprocess_type, is_still_mode, enhancer, batch_size, size_of_image, pose_style):
    """
    Create a lecture video combining slides and teacher video
    """
    try:
        if not slides_data:
            return None, "❌ Không có slide nào để xử lý!"
        
        # Create output directory
        output_dir = os.path.join("results", f"lecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}")
        os.makedirs(output_dir, exist_ok=True)
        
        print(f"Creating lecture video in: {output_dir}")
        
        # Copy source image to safe location
        safe_image_path = os.path.join(output_dir, "source_image.png")
        if not source_image or not os.path.exists(source_image):
            return None, "❌ Không tìm thấy ảnh nguồn!"
        
        shutil.copy2(source_image, safe_image_path)
        print(f"✅ Source image copied to safe location: {safe_image_path}")
        
        # Process each slide
        slide_clips = []
        total_duration = 0
        temp_video_files = []  # Track temporary video files for cleanup
        
        for i, slide_data in enumerate(slides_data):
            print(f"\n--- Processing slide {i+1}/{len(slides_data)} ---")
            
            # Verify source image exists before each slide processing
            if not os.path.exists(safe_image_path):
                print(f"⚠️ Source image lost, copying again...")
                if os.path.exists(source_image):
                    shutil.copy2(source_image, safe_image_path)
                    print(f"✅ Source image re-copied")
                else:
                    print(f"❌ Original source image also lost, stopping process")
                    break
            
            # Create slide image
            slide_image_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            slide_image = create_slide_image_with_text(slide_data['text'], slide_image_path)
            
            if not slide_image:
                print(f"❌ Failed to create slide image for slide {i+1}")
                continue
            
            # Generate audio for slide text
            audio_path = convert_text_to_audio(slide_data['text'], language)
            if not audio_path:
                print(f"❌ Failed to generate audio for slide {i+1}")
                continue
            
            # Get audio duration
            audio_duration = get_audio_duration(audio_path)
            print(f"Audio duration for slide {i+1}: {audio_duration:.2f} seconds")
            
            # If audio duration is 0 or very short, set a minimum duration
            if audio_duration <= 0.1:
                audio_duration = 3.0  # Minimum 3 seconds per slide
                print(f"⚠️ Audio duration too short, setting to minimum: {audio_duration}s")
            
            # Generate teacher video
            teacher_video_path = generate_video_for_text(
                sad_talker, safe_image_path, slide_data['text'], language, 
                preprocess_type, is_still_mode, enhancer, batch_size, size_of_image, pose_style
            )
            
            if not teacher_video_path or not os.path.exists(teacher_video_path):
                print(f"❌ Failed to generate teacher video for slide {i+1}")
                # Clean up audio
                if os.path.exists(audio_path):
                    os.remove(audio_path)
                continue
            
            # Add small delay to ensure file is fully written
            import time
            time.sleep(1)
            
            # Create slide video clip (static image for audio duration)
            slide_clip = ImageClip(slide_image_path, duration=audio_duration)
            
            # Load teacher video clip with error handling
            try:
                teacher_clip = VideoFileClip(teacher_video_path)
            except Exception as e:
                print(f"❌ Error loading teacher video for slide {i+1}: {str(e)}")
                # Clean up audio
                if os.path.exists(audio_path):
                    os.remove(audio_path)
                continue
            
            # Resize teacher video to fit in bottom-right corner (picture-in-picture)
            # Calculate size: 25% of slide width, maintain aspect ratio
            slide_width, slide_height = slide_clip.size
            teacher_width = int(slide_width * 0.25)
            teacher_height = int(teacher_width * teacher_clip.h / teacher_clip.w)
            
            # Resize teacher video
            teacher_clip = teacher_clip.resize((teacher_width, teacher_height))
            
            # Position teacher video in bottom-right corner
            teacher_x = slide_width - teacher_width - 50  # 50px margin
            teacher_y = slide_height - teacher_height - 50  # 50px margin
            teacher_clip = teacher_clip.set_position((teacher_x, teacher_y))
            
            # Composite slide and teacher video
            composite_clip = CompositeVideoClip([slide_clip, teacher_clip])
            
            # Add to clips list
            slide_clips.append(composite_clip)
            total_duration += audio_duration
            
            # Track video file for later cleanup
            temp_video_files.append(teacher_video_path)
            
            print(f"✅ Slide {i+1} processed: {audio_duration:.2f}s")
            
            # Clean up temporary files with delay
            time.sleep(0.5)  # Small delay before cleanup
            if os.path.exists(audio_path):
                try:
                    os.remove(audio_path)
                except Exception as e:
                    print(f"⚠️ Could not remove audio file: {str(e)}")
            
            # Don't delete teacher video file immediately - it might still be in use
            # We'll clean it up later in the final cleanup
            print(f"📁 Teacher video saved for slide {i+1}: {os.path.basename(teacher_video_path)}")
        
        if not slide_clips:
            return None, "❌ Không thể tạo video cho bất kỳ slide nào!"
        
        print(f"\n--- Creating final lecture video ---")
        print(f"Total slides: {len(slide_clips)}")
        print(f"Total duration: {total_duration:.2f} seconds")
        
        # Concatenate all slide clips
        final_video_path = os.path.join(output_dir, "lecture_final.mp4")
        final_video = concatenate_videoclips(slide_clips, method="compose")
        
        print(f"Writing final video to: {final_video_path}")
        final_video.write_videofile(final_video_path, codec='libx264', audio_codec='aac', verbose=False, logger=None)
        
        # Clean up clips
        for clip in slide_clips:
            try:
                clip.close()
            except Exception as e:
                print(f"⚠️ Could not close clip: {str(e)}")
        try:
            final_video.close()
        except Exception as e:
            print(f"⚠️ Could not close final video: {str(e)}")
        
        # Add delay before cleanup
        time.sleep(1)
        
        # Clean up temporary slide images
        for i in range(len(slides_data)):
            slide_image_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            if os.path.exists(slide_image_path):
                try:
                    os.remove(slide_image_path)
                except Exception as e:
                    print(f"⚠️ Could not remove slide image {i+1}: {str(e)}")
        
        # Clean up source image copy
        if os.path.exists(safe_image_path):
            try:
                os.remove(safe_image_path)
            except Exception as e:
                print(f"⚠️ Could not remove source image copy: {str(e)}")
        
        # Clean up temporary video files
        print("🧹 Cleaning up temporary video files...")
        for video_file in temp_video_files:
            if os.path.exists(video_file):
                try:
                    os.remove(video_file)
                    print(f"  ✅ Deleted: {os.path.basename(video_file)}")
                except Exception as e:
                    print(f"  ❌ Could not delete {os.path.basename(video_file)}: {str(e)}")
        
        # Also clean up temporary SadTalker directories
        print("🗂️ Cleaning up temporary SadTalker directories...")
        temp_dirs_deleted = 0
        results_dir = "results"
        if os.path.exists(results_dir):
            for item in os.listdir(results_dir):
                item_path = os.path.join(results_dir, item)
                # Check if it's a directory with UUID-like name (temporary SadTalker dirs)
                if os.path.isdir(item_path) and len(item) == 36 and '-' in item:
                    try:
                        shutil.rmtree(item_path)
                        print(f"  ✅ Deleted temp directory: {item}")
                        temp_dirs_deleted += 1
                    except Exception as e:
                        print(f"  ❌ Failed to delete temp directory {item}: {str(e)}")
        
        print(f"🗂️ Cleanup completed: {temp_dirs_deleted} temporary directories deleted")
        
        print(f"✅ Lecture video created successfully: {final_video_path}")
        
        return final_video_path, f"✅ Hoàn thành! Đã tạo video bài giảng với {len(slides_data)} slide, tổng thời gian: {total_duration:.1f}s"
        
    except Exception as e:
        print(f"Error in create_lecture_video: {str(e)}")
        return None, f"❌ Lỗi tạo video bài giảng: {str(e)}"

def generate_video_for_text(sad_talker, source_image, text, language, preprocess_type, is_still_mode, enhancer, batch_size, size_of_image, pose_style):
    """
    Generate video for a single text using SadTalker
    """
    try:
        # Verify source image exists
        if not source_image or not os.path.exists(source_image):
            print(f"Source image not found: {source_image}")
            return None
        
        # Convert text to audio
        audio_path = convert_text_to_audio(text, language)
        if not audio_path:
            print("Failed to convert text to audio")
            return None
        
        print(f"Generating video with image: {source_image}")
        print(f"Image exists: {os.path.exists(source_image)}")
        print(f"Audio path: {audio_path}")
        
        # Generate video using the same method as the main interface
        video_path = sad_talker.test(
            source_image, audio_path, preprocess_type, is_still_mode, 
            enhancer, batch_size, size_of_image, pose_style
        )
        
        # Clean up temporary audio file
        if os.path.exists(audio_path):
            os.remove(audio_path)
        
        if video_path and os.path.exists(video_path):
            print(f"✅ Video generated successfully: {video_path}")
            return video_path
        else:
            print(f"❌ Video generation failed or file not found: {video_path}")
            return None
            
    except Exception as e:
        print(f"Error generating video for text: {str(e)}")
        return None

def merge_videos_with_transition(video_paths, output_path, transition_duration=1.0):
    """
    Merge multiple videos with simple fade transition
    """
    try:
        if not video_paths:
            print("No video paths provided for merging")
            return None
        
        print(f"Merging {len(video_paths)} videos...")
        
        clips = []
        for i, video_path in enumerate(video_paths):
            if os.path.exists(video_path):
                print(f"Loading video {i+1}: {video_path}")
                clip = VideoFileClip(video_path)
                
                # Add fade in/out effects
                if i == 0:  # First clip
                    clip = clip.fadein(transition_duration)
                if i == len(video_paths) - 1:  # Last clip
                    clip = clip.fadeout(transition_duration)
                else:  # Middle clips
                    clip = clip.fadein(transition_duration/2).fadeout(transition_duration/2)
                
                clips.append(clip)
                print(f"✅ Video {i+1} loaded and processed")
            else:
                print(f"❌ Video file not found: {video_path}")
        
        if not clips:
            print("No valid video clips to merge")
            return None
        
        print(f"Concatenating {len(clips)} clips...")
        
        # Concatenate all clips
        final_video = concatenate_videoclips(clips, method="compose")
        
        print(f"Writing final video to: {output_path}")
        final_video.write_videofile(output_path, codec='libx264', audio_codec='aac', verbose=False, logger=None)
        
        # Clean up clips
        for clip in clips:
            clip.close()
        final_video.close()
        
        print(f"✅ Video merge completed: {output_path}")
        return output_path
    except Exception as e:
        print(f"Error in merge_videos_with_transition: {str(e)}")
        return None

def process_powerpoint_to_video(sad_talker, pptx_file, source_image, language, preprocess_type, is_still_mode, enhancer, batch_size, size_of_image, pose_style):
    """
    Main function to process PowerPoint file and generate final video
    """
    try:
        # Extract text from PowerPoint
        slides_text = extract_text_from_pptx(pptx_file)
        if not slides_text:
            return None, "❌ Không thể đọc nội dung từ file PowerPoint!"
        
        print(f"Extracted {len(slides_text)} slides from PowerPoint")
        for i, text in enumerate(slides_text):
            print(f"Slide {i+1}: {text[:50]}...")
        
        # Create output directory based on PowerPoint filename
        pptx_name = os.path.splitext(os.path.basename(pptx_file.name))[0]
        output_dir = os.path.join("results", f"powerpoint_{pptx_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}")
        os.makedirs(output_dir, exist_ok=True)
        
        print(f"Output directory: {output_dir}")
        
        # Copy source image to safe location first
        safe_image_path = os.path.join(output_dir, "source_image.png")
        if not source_image or not os.path.exists(source_image):
            return None, "❌ Không tìm thấy ảnh nguồn!"
        
        # Copy image to safe location
        shutil.copy2(source_image, safe_image_path)
        print(f"✅ Source image copied to safe location: {safe_image_path}")
        
        # Ensure the output directory exists and is writable
        if not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        print(f"Output directory verified: {output_dir}")
        
        # Pre-warm the model for better performance (optional)
        # Uncomment the lines below if you want to enable pre-warming
        # print("🔥 Pre-warming SadTalker model...")
        # try:
        #     # Generate a short test video to warm up the model
        #     test_audio_path = convert_text_to_audio("Test", language)
        #     if test_audio_path and os.path.exists(test_audio_path):
        #         sad_talker.test(
        #             safe_image_path, test_audio_path, preprocess_type, is_still_mode, 
        #             enhancer, batch_size, size_of_image, pose_style
        #         )
        #         # Clean up test audio
        #         if os.path.exists(test_audio_path):
        #             os.remove(test_audio_path)
        #         print("✅ Model pre-warmed successfully")
        # except Exception as e:
        #     print(f"⚠️ Model pre-warming failed: {str(e)}")
        
        print("🚀 Starting video generation without pre-warming...")
        
        # Generate videos for each slide using loop
        video_paths = []
        total_slides = len(slides_text)
        
        print(f"Starting video generation for {total_slides} slides...")
        
        for i, text in enumerate(slides_text):
            if text.strip():
                print(f"\n--- Processing slide {i+1}/{total_slides} ---")
                print(f"Text content: {text[:100]}...")
                
                # Verify safe image exists before each generation
                if not os.path.exists(safe_image_path):
                    print(f"❌ Safe image lost during processing: {safe_image_path}")
                    # Try to copy again
                    if os.path.exists(source_image):
                        shutil.copy2(source_image, safe_image_path)
                        print(f"✅ Image re-copied to safe location")
                    else:
                        print(f"❌ Original image also lost, stopping process")
                        break
                
                # Generate video for this slide
                print(f"Generating video for slide {i+1}...")
                print(f"Using safe image path: {safe_image_path}")
                print(f"Safe image exists: {os.path.exists(safe_image_path)}")
                video_path = generate_video_for_text(
                    sad_talker, safe_image_path, text, language, preprocess_type, 
                    is_still_mode, enhancer, batch_size, size_of_image, pose_style
                )
                
                if video_path and os.path.exists(video_path):
                    # Save video to output directory
                    slide_video_path = os.path.join(output_dir, f"slide_{i+1:02d}.mp4")
                    shutil.move(video_path, slide_video_path)
                    video_paths.append(slide_video_path)
                    print(f"✅ Slide {i+1} video saved: {slide_video_path}")
                else:
                    print(f"❌ Failed to generate video for slide {i+1}")
                    continue
            else:
                print(f"Skipping slide {i+1} - empty text")
        
        print(f"\n--- Video Generation Summary ---")
        print(f"Total slides processed: {total_slides}")
        print(f"Videos successfully generated: {len(video_paths)}")
        for i, path in enumerate(video_paths):
            print(f"  Video {i+1}: {path}")
        
        if not video_paths:
            return None, "❌ Không thể tạo video cho bất kỳ slide nào!"
        
        print(f"Generated {len(video_paths)} videos. Starting merge...")
        
        # Merge all videos
        final_video_path = os.path.join(output_dir, f"{pptx_name}_final.mp4")
        merged_video = merge_videos_with_transition(video_paths, final_video_path)
        
        if merged_video and os.path.exists(final_video_path):
            print(f"✅ Final video created: {final_video_path}")
            
            # Clean up individual slide videos to save space
            print("🧹 Cleaning up individual slide videos...")
            deleted_count = 0
            for video_path in video_paths:
                if os.path.exists(video_path):
                    try:
                        os.remove(video_path)
                        print(f"  ✅ Deleted: {os.path.basename(video_path)}")
                        deleted_count += 1
                    except Exception as e:
                        print(f"  ❌ Failed to delete {os.path.basename(video_path)}: {str(e)}")
            
            # Also clean up the source image copy
            if os.path.exists(safe_image_path):
                try:
                    os.remove(safe_image_path)
                    print(f"  ✅ Deleted: source_image.png")
                except Exception as e:
                    print(f"  ❌ Failed to delete source_image.png: {str(e)}")
            
            print(f"🧹 Cleanup completed: {deleted_count}/{len(video_paths)} slide videos deleted")
            
            # Clean up temporary SadTalker directories
            print("🗂️ Cleaning up temporary SadTalker directories...")
            temp_dirs_deleted = 0
            results_dir = "results"
            if os.path.exists(results_dir):
                for item in os.listdir(results_dir):
                    item_path = os.path.join(results_dir, item)
                    # Check if it's a directory with UUID-like name (temporary SadTalker dirs)
                    if os.path.isdir(item_path) and len(item) == 36 and '-' in item:
                        try:
                            shutil.rmtree(item_path)
                            print(f"  ✅ Deleted temp directory: {item}")
                            temp_dirs_deleted += 1
                        except Exception as e:
                            print(f"  ❌ Failed to delete temp directory {item}: {str(e)}")
            
            print(f"🗂️ Cleanup completed: {temp_dirs_deleted} temporary directories deleted")
            
            # Calculate space saved
            total_size_saved = 0
            for video_path in video_paths:
                if os.path.exists(video_path):
                    total_size_saved += os.path.getsize(video_path)
            
            if total_size_saved > 0:
                size_mb = total_size_saved / (1024 * 1024)
                print(f"💾 Space saved: {size_mb:.2f} MB")
            
            cleanup_summary = f"✅ Hoàn thành! Đã tạo {len(video_paths)} video, gộp thành 1 video cuối cùng"
            if deleted_count > 0:
                cleanup_summary += f", dọn dẹp {deleted_count} file slide"
            if temp_dirs_deleted > 0:
                cleanup_summary += f" và {temp_dirs_deleted} thư mục tạm"
            cleanup_summary += f". Lưu trong: {output_dir}"
            
            return final_video_path, cleanup_summary
        else:
            return None, "❌ Lỗi khi gộp video!"
            
    except Exception as e:
        print(f"Error in process_powerpoint_to_video: {str(e)}")
        return None, f"❌ Lỗi xử lý PowerPoint: {str(e)}"

# CSS nhẹ chỉ cho home (card bo góc, shadow, xanh nhạt)
custom_home_css = """
.sad-home-cards-row {
    display: flex; gap: 52px; justify-content: center; margin: 52px 0 36px 0;
}
.sad-home-card {
    background: #f5f9ff;
    border-radius: 30px;
    box-shadow: 0 12px 40px rgba(50,130,255,0.16), 0 2px 10px rgba(180,192,240,0.13);
    padding: 44px 54px;
    min-width: 350px; max-width: 540px; width: 100%;
    border: 2.8px solid #b1d1ff;
    transition: box-shadow 0.2s, border 0.2s;
}
.sad-home-card:hover {
    box-shadow: 0 16px 60px rgba(44, 110, 255,0.22), 0 4px 16px rgba(60,110,255,0.17);
    border: 2.8px solid #347cff;
}
.sad-home-card h3 {
    color: #1460a5;
    font-size: 2rem;
    font-weight: 900;
    margin-bottom: 24px;
    display: flex; align-items: center; gap: 13px;
    letter-spacing: 0.5px;
    text-shadow: 0 2px 8px #e0edfa;
}
.sad-home-card ul, .sad-home-card ol {
    font-size: 1.19rem;
    color: #000;
    font-weight: 600;
    margin-left: 12px; margin-top: 3px; line-height: 1.85;
    letter-spacing: 0.2px;
}
.sad-home-card li {
margin-bottom: 10px;
list-style:none;
}
.sad-home-btn-wrap {display: flex; justify-content: center; margin: 18px 0 6px 0;}
.sad-home-btn {
    background: linear-gradient(90deg, #468fff, #54b1ff 80%);
    color: #fff; border: none; border-radius: 16px; font-size: 1.18rem;
    padding: 18px 50px; font-weight: bold; letter-spacing: 1px;
    box-shadow: 0 6px 18px rgba(48,124,247,0.12);
    cursor: pointer; transition: background .13s, transform .12s;
}
.sad-home-btn:hover {
    background: linear-gradient(90deg, #357aff, #468fff 80%);
    transform: translateY(-2px) scale(1.04);
}
.sad-home-title {
    text-align: center; color: #2165c3; font-size: 2.5rem;
    font-weight: bold; margin-top: 38px; margin-bottom: 12px; letter-spacing: 1.1px;
}
.sad-home-desc {text-align: center; color: #222; font-size: 1.23rem; opacity: 0.92;}
@media (max-width: 1100px){
    .sad-home-cards-row { flex-direction: column; align-items: center; gap: 26px;}
    .sad-home-card { min-width: 200px; max-width: 99vw;}
}
"""


def sadtalker_demo_with_home(checkpoint_path='checkpoints', config_path='src/config', warpfn=None):
    sad_talker = SadTalker(checkpoint_path, config_path, lazy_load=True)

    with gr.Blocks(analytics_enabled=False, title="SadTalker", css=custom_home_css) as sadtalker_interface:
        with gr.Tabs():
            # --- TAB HOME ---
            with gr.TabItem("🏠 Trang chủ"):
                gr.HTML("""
                <div class="sad-home-title"><span style="font-size:2.2rem;">🎭</span> SadTalker</div>
                <div class="sad-home-desc">Tạo video nói chuyện từ ảnh tĩnh và âm thanh với AI<br>
                    <span style="font-size:0.98rem;">Chào mừng bạn đến với SadTalker - Công cụ tạo video nói chuyện thông minh dựa trên nghiên cứu CVPR 2023</span>
                </div>
                <div class="sad-home-cards-row">
                    <div class="sad-home-card">
                        <h3 style="color: #000 !important;">🚀 Tính năng nổi bật</h3>
                        <ul style="color: #000 !important; font-size: 1.19rem; font-weight: 600; margin-left: 12px; margin-top: 3px; line-height: 1.85; letter-spacing: 0.2px;">
                            <li style="color: #000 !important; margin-bottom: 10px;">✨ Tạo video nói chuyện từ ảnh tĩnh</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🎵 Hỗ trợ âm thanh từ file hoặc văn bản</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">📄 Import văn bản từ file (.txt, .md, .doc, .docx)</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🎓 Tạo video bài giảng (slide + giáo viên giảng)</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🎨 Nhiều tùy chọn xử lý ảnh</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">⚡ Tốc độ xử lý nhanh</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🤖 Hiệu ứng chuyển động tự nhiên</li>
                        </ul>
                    </div>
                    <div class="sad-home-card">
                        <h3 style="color: #000 !important;">💡 Hướng dẫn sử dụng</h3>
                        <ul style="color: #000 !important; font-size: 1.19rem; font-weight: 600; margin-left: 12px; margin-top: 3px; line-height: 1.85; letter-spacing: 0.2px;">
                            <li style="color: #000 !important; margin-bottom: 10px;">📸 Tải lên ảnh khuôn mặt</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🎤 Chọn file âm thanh hoặc nhập văn bản</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">📄 Hoặc import văn bản từ file</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">⚙️ Điều chỉnh các thông số</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">🎬 Nhấn "Sinh Video" để tạo</li>
                            <li style="color: #000 !important; margin-bottom: 10px;">💾 Tải xuống kết quả</li>
                        </ul>
                    </div>
                </div>
                <div class="sad-home-btn-wrap">
                    <button class="sad-home-btn" onclick="document.querySelectorAll('.tabitem')[1].click();">🎓 Chuyển đến giao diện Tạo Bài Giảng</button>
                </div>
                """)
            

            
            # --- TAB TẠO BÀI GIẢNG ---
            with gr.TabItem("🎓 Tạo Bài Giảng"):
                gr.Markdown("<div align='center'> <h2> 🎓 SadTalker Lecture Video Generator </h2> \
                    <p>Tạo video bài giảng kết hợp slide PowerPoint và video giáo viên giảng bài</p> </div>")

                with gr.Row().style(equal_height=False):
                    with gr.Column(variant='panel'):
                        # Source image upload
                        gr.Markdown("### 👨‍🏫 Ảnh Giáo Viên")
                        lecture_source_image = gr.Image(label="Ảnh khuôn mặt giáo viên", source="upload", type="filepath", elem_id="lecture_source_image").style(width=512)
                        
                        # PowerPoint file upload
                        gr.Markdown("### 📊 File PowerPoint Bài Giảng")
                        lecture_pptx_file = gr.File(
                            label="Chọn file PowerPoint (.pptx)",
                            file_types=[".pptx"],
                            elem_id="lecture_pptx_file"
                        )
                        
                        # Language selection
                        lecture_audio_language = gr.Dropdown(
                            choices=["vi", "en", "zh", "ja", "ko", "fr", "de", "es", "it", "pt"],
                            value="vi",
                            label="Ngôn ngữ giảng bài",
                            elem_id="lecture_audio_language"
                        )
                        
                        # Preview extracted slides
                        gr.Markdown("### 📝 Nội dung từ PowerPoint")
                        lecture_slides_preview = gr.Textbox(
                            label="Nội dung đã trích xuất từ các slide",
                            lines=8,
                            interactive=False,
                            elem_id="lecture_slides_preview"
                        )
                        
                        # Extract slides button
                        extract_lecture_slides_btn = gr.Button(
                            '📂 Trích xuất nội dung từ PowerPoint',
                            elem_id="extract_lecture_slides_btn",
                            variant='secondary'
                        )
                        
                        # Generate lecture video button
                        generate_lecture_btn = gr.Button(
                            '🎬 Tạo Video Bài Giảng',
                            elem_id="generate_lecture_btn",
                            variant='primary'
                        )
                        
                        # Status
                        lecture_status = gr.Textbox(
                            label="Trạng thái xử lý",
                            interactive=False,
                            elem_id="lecture_status"
                        )
                    
                    with gr.Column(variant='panel'):
                        # Settings
                        with gr.Tabs(elem_id="lecture_settings"):
                            with gr.TabItem('⚙️ Cài đặt'):
                                gr.Markdown("Cài đặt cho video bài giảng")
                                with gr.Column(variant='panel'):
                                    lecture_pose_style = gr.Slider(minimum=0, maximum=46, step=1, label="Pose style", value=0)
                                    lecture_size_of_image = gr.Radio([256, 512], value=256, label='Độ phân giải ảnh', info="256 = Nhanh, 512 = Chất lượng cao")
                                    lecture_preprocess_type = gr.Radio(['crop', 'resize','full', 'extcrop', 'extfull'], value='crop', label='Xử lý ảnh', info="crop = Nhanh nhất")
                                    lecture_is_still_mode = gr.Checkbox(label="Still Mode (ít chuyển động đầu)")
                                    lecture_batch_size = gr.Slider(label="Batch size", step=1, maximum=10, value=6, info="Tăng lên 6-8 để nhanh hơn")
                                    lecture_enhancer = gr.Checkbox(label="GFPGAN làm Face enhancer (chậm hơn)")
                                    
                                    # Fast mode preset button
                                    lecture_fast_mode_btn = gr.Button(
                                        '⚡ Chế độ nhanh (256px, batch=6, không enhancer)',
                                        elem_id="lecture_fast_mode_btn",
                                        variant='secondary'
                                    )
                                    
                                    def set_lecture_fast_mode():
                                        return 256, 'crop', False, 6, False, 0
                                    
                                    lecture_fast_mode_btn.click(
                                        fn=set_lecture_fast_mode,
                                        outputs=[lecture_size_of_image, lecture_preprocess_type, lecture_is_still_mode, lecture_batch_size, lecture_enhancer, lecture_pose_style]
                                    )
                        
                        # Results
                        with gr.Tabs(elem_id="lecture_results"):
                            gr.Markdown("### 🎬 Video Bài Giảng")
                            lecture_final_video = gr.Video(label="Video bài giảng hoàn chỉnh", format="mp4").style(width=512)
                            
                            gr.Markdown("### 📊 Thông tin Video")
                            lecture_info = gr.Textbox(
                                label="Thông tin chi tiết",
                                lines=4,
                                interactive=False,
                                elem_id="lecture_info"
                            )
                
                # Lecture event handlers
                def extract_lecture_slides(file):
                    if file:
                        slides_data = extract_slides_from_pptx(file)
                        if slides_data:
                            slides_text = []
                            for i, slide in enumerate(slides_data):
                                slides_text.append(f"Slide {i+1}: {slide['text'][:100]}...")
                            return "\n\n".join(slides_text)
                    return "❌ Vui lòng chọn file PowerPoint!"
                
                def generate_lecture_video_handler(pptx, img, lang, preprocess, still, enh, batch, size, pose):
                    if not pptx or not img:
                        return None, "❌ Vui lòng chọn file PowerPoint và ảnh giáo viên!"
                    
                    slides_data = extract_slides_from_pptx(pptx)
                    if not slides_data:
                        return None, "❌ Không thể trích xuất nội dung từ PowerPoint!"
                    
                    return create_lecture_video(
                        sad_talker, slides_data, img, lang, preprocess, still, enh, batch, size, pose
                    )
                
                extract_lecture_slides_btn.click(
                    fn=extract_lecture_slides,
                    inputs=[lecture_pptx_file],
                    outputs=[lecture_slides_preview]
                )
                
                generate_lecture_btn.click(
                    fn=generate_lecture_video_handler,
                    inputs=[
                        lecture_pptx_file, lecture_source_image, lecture_audio_language,
                        lecture_preprocess_type, lecture_is_still_mode, lecture_enhancer, lecture_batch_size, lecture_size_of_image, lecture_pose_style
                    ],
                    outputs=[lecture_final_video, lecture_status]
                )
    return sadtalker_interface

if __name__ == "__main__":
    demo = sadtalker_demo_with_home()
    demo.launch(
        server_name='127.0.0.1',
        server_port=7860,
        show_error=True,
        quiet=False,
        share=False,
        inbrowser=True
    )