import os
import zipfile
import subprocess
import tempfile
import logging
import gradio as gr
import xml.etree.ElementTree as ET
from shutil import which
from gtts import gTTS
from pptx import Presentation
from pdf2image import convert_from_path
from src.utils.math_formula_processor import MathFormulaProcessor, process_math_text
from src.voice.enrollment import enroll_voice
# Thêm import backend voice
from src.voice.store import list_voices, has_voice
from src.voice.tts_engine import synthesize

# Thiết lập logging
logger = logging.getLogger(__name__)


def convert_pptx_to_images(pptx_path, dpi=220):
    if which("soffice") is None:
        raise RuntimeError("Không tìm thấy LibreOffice (soffice). Vui lòng cài LibreOffice để chuyển PPTX -> PDF.")

    tmpdir = tempfile.mkdtemp(prefix="pptx2img_")
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Chuyển PPTX sang PDF thất bại: {e}")

    pdf_path = os.path.join(tmpdir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(pdf_path):
        raise RuntimeError("Không tạo được PDF từ PPTX. Kiểm tra file đầu vào.")

    try:
        images = convert_from_path(pdf_path, dpi=dpi, output_folder=tmpdir, fmt='png')
    except Exception as e:
        raise RuntimeError("Lỗi convert PDF -> ảnh. Có thể thiếu Poppler (poppler-utils).") from e

    img_paths = []
    for i, img in enumerate(images, 1):
        img_path = os.path.join(tmpdir, f"slide-{i:02d}.png")
        img.save(img_path)
        img_paths.append(img_path)
    return img_paths


def convert_text_to_audio(text, language='vi'):
    try:
        if not text or text.strip() == "":
            return None
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
        temp_path = temp_file.name
        temp_file.close()
        tts = gTTS(text=text, lang=language, slow=False)
        tts.save(temp_path)
        return temp_path
    except Exception:
        return None

def convert_text_to_audio_with_voice(text, user_id, voice_id, language=None):
    """
    Convert text to audio using registered voice (XTTS v2)
    Falls back to gTTS if XTTS v2 is not available
    """
    try:
        if not text or text.strip() == "":
            return None
            
        print(f"🔍 Debug: text='{text[:50]}...', user_id='{user_id}', voice_id='{voice_id}'")
            
        # Bắt buộc phải có voice
        if not has_voice(user_id, voice_id):
            print(f"❌ Voice not found: {voice_id}")
            raise Exception(f"Giọng nói '{voice_id}' không tồn tại. Vui lòng đăng ký giọng nói trước!")
        
        print(f"🎤 Using registered voice: {voice_id}")
        
        # Lấy thông tin voice để biết ngôn ngữ
        from src.voice.store import load_profile
        emb, meta = load_profile(user_id, voice_id)
        voice_language = meta.get('lang_hint', 'vi')
        model_type = meta.get('model_type', 'unknown')
        
        print(f"🌐 Using voice language: {voice_language}")
        print(f"🤖 Model type: {model_type}")
        
        # Kiểm tra nếu model là Tacotron2 (không hỗ trợ voice cloning)
        if model_type == "tacotron2":
            print(f"⚠️ Voice was registered with Tacotron2 (no voice cloning support)")
            print(f"🔊 Falling back to gTTS with voice language: {voice_language}")
            return convert_text_to_audio(text, voice_language)
        
        # Sử dụng giọng nói đã đăng ký với ngôn ngữ của voice
        audio_path = synthesize(text, user_id, voice_id, voice_language)
        print(f"🔍 synthesize() returned: {audio_path}")
        
        if audio_path and os.path.exists(audio_path):
            print(f"✅ Successfully synthesized with registered voice: {audio_path}")
            return audio_path
        else:
            print(f"❌ Failed to synthesize with registered voice")
            print(f"🔍 audio_path: {audio_path}")
            print(f"🔍 exists: {os.path.exists(audio_path) if audio_path else 'None'}")
            print(f"🔊 Falling back to gTTS with voice language: {voice_language}")
            return convert_text_to_audio(text, voice_language)
        
    except Exception as e:
        print(f"❌ Error in convert_text_to_audio_with_voice: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Fallback to gTTS with voice language
        try:
            voice_language = 'vi'  # Default fallback
            if has_voice(user_id, voice_id):
                emb, meta = load_profile(user_id, voice_id)
                voice_language = meta.get('lang_hint', 'vi')
            
            print(f"🔊 Final fallback to gTTS with language: {voice_language}")
            return convert_text_to_audio(text, voice_language)
        except Exception as fallback_error:
            print(f"❌ Even gTTS fallback failed: {fallback_error}")
            raise e


def read_text_file(file):
    if file is None:
        return "", "❌ Vui lòng chọn file văn bản!"
    try:
        file_path = file.name
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in ['.txt', '.md', '.doc', '.docx', '.rtf']:
            return "", "❌ Chỉ hỗ trợ file văn bản (.txt, .md, .doc, .docx, .rtf)"
        content = ""
        if file_ext == '.docx':
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    if 'word/document.xml' in zip_file.namelist():
                        xml_content = zip_file.read('word/document.xml')
                        root = ET.fromstring(xml_content)
                        text_elements = []
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                text_elements.append(elem.text.strip())
                        content = ' '.join(text_elements)
                    else:
                        return "", "❌ Không thể đọc nội dung từ file .docx"
            except Exception as e:
                return "", f"❌ Lỗi đọc file .docx: {str(e)}"
        elif file_ext == '.doc':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        content = f.read()
                except:
                    return "", "❌ Không thể đọc file .doc. Vui lòng chuyển đổi sang .txt hoặc .docx"
        else:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
        content = ''.join(char for char in content if ord(char) >= 32 or char in '\n\r\t').strip()
        if not content:
            return "", "❌ File văn bản trống hoặc chỉ chứa ký tự đặc biệt!"
        return content, f"✅ Đã đọc thành công file: {os.path.basename(file_path)} ({len(content)} ký tự thực tế)"
    except Exception as e:
        return "", f"❌ Lỗi đọc file: {str(e)}"


def extract_slides_from_pptx(pptx_file):
    slides_data = []
    image_paths = convert_pptx_to_images(pptx_file.name, dpi=220)
    math_processor = MathFormulaProcessor()
    processed_result = math_processor.process_powerpoint_text(pptx_file.name)

    if processed_result.get('error'):
        logger.warning(f"Lỗi xử lý công thức toán học: {processed_result['error']}, sử dụng phương pháp cũ")
        prs = Presentation(pptx_file.name)
        for i, slide in enumerate(prs.slides):
            text_chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_chunks.append(shape.text.strip())
            text = " ".join(filter(None, text_chunks))
            processed_text = process_math_text(text.strip())
            slides_data.append({
                'slide_number': i + 1,
                'text': processed_text,
                'image_path': image_paths[i] if i < len(image_paths) else None,
                'has_math_objects': False
            })
    else:
        for slide_info in processed_result['slides']:
            slides_data.append({
                'slide_number': slide_info['slide_number'],
                'text': slide_info['processed_text'],
                'image_path': image_paths[slide_info['slide_number'] - 1] if slide_info['slide_number'] - 1 < len(image_paths) else None,
                'has_math_objects': slide_info['has_math_objects']
            })
    return slides_data


def extract_lecture_slides(file):
    if file:
        slides_data = extract_slides_from_pptx(file)
        if slides_data:
            slides_text = []
            total_math_slides = 0
            for slide in slides_data:
                slide_num = slide['slide_number']
                text_preview = slide['text'][:100] + "..." if len(slide['text']) > 100 else slide['text']
                if slide.get('has_math_objects', False):
                    slides_text.append(f"Slide {slide_num}: {text_preview} [📐 Có công thức toán học]")
                    total_math_slides += 1
                else:
                    slides_text.append(f"Slide {slide_num}: {text_preview}")
            summary = f"\n\n📊 Thống kê: {len(slides_data)} slides, {total_math_slides} slides có công thức toán học"
            if total_math_slides > 0:
                summary += "\n✅ Đã xử lý thành công các ký tự đặc biệt và công thức toán học!"
            return "\n\n".join(slides_text) + summary
    return "❌ Vui lòng chọn file PowerPoint!"


def set_lecture_fast_mode():
    return 256, 'crop', False, 6, False, 0

def handle_register_voice(file_obj):
    if file_obj is None:
        return "❌ Chưa upload file", gr.update(choices=[])
    
    audio_path = file_obj.name  # lấy đường dẫn tạm mà Gradio lưu file
    voice_id = os.path.basename(audio_path).split('.')[0]
    
    print(f"🔍 Debug: Đang đăng ký voice '{voice_id}' từ file: {audio_path}")
    
    # Kiểm tra voice đã tồn tại chưa
    if has_voice("current_user", voice_id):
        print(f"⚠️  Voice '{voice_id}' đã tồn tại, sẽ ghi đè")
    
    # Sử dụng ngôn ngữ mặc định là 'vi' cho giọng nhân bản
    success = enroll_voice(audio_path, "current_user", voice_id, lang_hint="vi")
    
    if not success:
        print(f"❌ Đăng ký voice '{voice_id}' thất bại")
        return "❌ Đăng ký giọng thất bại", gr.update(choices=[])
    
    print(f"✅ Đăng ký voice '{voice_id}' thành công")
    
    # Đợi một chút để đảm bảo file được lưu
    import time
    time.sleep(0.5)
    
    # Refresh dropdown với danh sách giọng mới
    voices = list_voices("current_user")
    voice_choices = [m["voice_id"] for m in voices]
    
    print(f"🔍 Debug: Danh sách voices sau khi đăng ký: {voice_choices}")
    print(f"🔍 Debug: Voice ID hiện tại: {voice_id}")
    print(f"🔍 Debug: Voice ID có trong danh sách: {voice_id in voice_choices}")
    
    # Tìm voice profile mới nhất để lấy thông tin
    if voices:
        latest_voice = voices[0]  # voices đã được sort theo created_at
        duration = latest_voice.get("audio_length", 0) / latest_voice.get("sample_rate", 48000)
        print(f"✅ Voice '{voice_id}' đã được lưu với duration: {duration:.1f}s")
        return f"✅ Đã lưu giọng '{voice_id}' (~{duration:.1f}s)", gr.update(choices=voice_choices, value=voice_id)
    else:
        print(f"⚠️  Không tìm thấy voices sau khi đăng ký")
        return f"✅ Đã lưu giọng '{voice_id}'", gr.update(choices=voice_choices, value=voice_id)

def create_lecture_input_interface():
    with gr.Row(equal_height=False):        
        with gr.Column(variant='panel'):
            gr.Markdown("### 👨‍🏫 Ảnh Giáo Viên")
            lecture_source_image = gr.Image(label="Ảnh khuôn mặt giáo viên", source="upload", type="filepath")

            gr.Markdown("### 📊 File PowerPoint Bài Giảng")
            lecture_pptx_file = gr.File(label="Chọn file PowerPoint (.pptx)", file_types=[".pptx"])

            gr.Markdown("### 🎤 Giọng Nhân Bản (Bắt Buộc)")
            gr.Markdown("*⚠️ Bạn phải đăng ký giọng nói trước khi tạo video bài giảng*")
            
            # Upload audio voice sample
            lecture_voice_sample = gr.File(
                label="📢 Upload audio mẫu để nhân bản giọng",
                file_types=[".wav", ".mp3"],
                type="file"  # thay filepath thành file
            )
            # Nút đăng ký giọng
            register_voice_btn = gr.Button("💾 Đăng ký giọng", variant="primary")
            register_voice_status = gr.Textbox(
                label="Trạng thái đăng ký giọng",
                interactive=False
            )
            
            # Button debug để kiểm tra voice store
            debug_voice_btn = gr.Button("🔍 Debug Voice Store", variant="secondary", size="sm")

            lecture_voice_id = gr.Dropdown(
                choices=[m["voice_id"] for m in list_voices("current_user")],
                label="🎤 Chọn giọng nhân bản (Bắt buộc)",
                value=None,
                info="Ngôn ngữ sẽ được lấy từ giọng đã đăng ký"
            )

            # Khi bấm nút -> gọi hàm enroll_voice và refresh dropdown
            register_voice_btn.click(
                fn=handle_register_voice,
                inputs=[lecture_voice_sample],
                outputs=[register_voice_status, lecture_voice_id]
            )
            
            # Thêm event để refresh dropdown khi có thay đổi
            def refresh_voice_dropdown():
                voices = list_voices("current_user")
                voice_choices = [m["voice_id"] for m in voices]
                print(f"🔍 Debug: Refreshing voice dropdown, found {len(voices)} voices: {voice_choices}")
                return gr.update(choices=voice_choices)

            def debug_voice_store():
                """Debug function để kiểm tra voice store"""
                print("🔍 Debug: Kiểm tra voice store...")
                
                # Kiểm tra thư mục voices
                from src.voice.store import ROOT
                import os
                voices_dir = ROOT / "current_user"
                print(f"🔍 Voice directory: {voices_dir}")
                print(f"🔍 Directory exists: {voices_dir.exists()}")
                
                if voices_dir.exists():
                    for item in voices_dir.iterdir():
                        print(f"🔍 Found item: {item.name} (dir: {item.is_dir()})")
                        if item.is_dir():
                            meta_file = item / "meta.json"
                            emb_file = item / "embedding.npy"
                            print(f"  - meta.json exists: {meta_file.exists()}")
                            print(f"  - embedding.npy exists: {emb_file.exists()}")
                
                # Kiểm tra list_voices function
                voices = list_voices("current_user")
                print(f"🔍 list_voices result: {voices}")
                
                return f"Debug complete. Found {len(voices)} voices."
            
            # Button debug voice store
            debug_voice_btn.click(
                fn=debug_voice_store,
                outputs=[register_voice_status]
            )
            
            # Refresh dropdown khi có thay đổi
            lecture_voice_sample.change(
                fn=refresh_voice_dropdown,
                outputs=[lecture_voice_id]
            )

            gr.Markdown("### 📝 Nội dung từ PowerPoint")
            lecture_slides_preview = gr.Textbox(label="Nội dung đã trích xuất từ các slide", lines=8, interactive=False)

            extract_lecture_slides_btn = gr.Button('📂 Trích xuất nội dung từ PowerPoint', variant='secondary')
            generate_lecture_btn = gr.Button('🎬 Tạo Video Bài Giảng', variant='primary')
            lecture_status = gr.Textbox(label="Trạng thái xử lý", interactive=False)

        with gr.Column(variant='panel'):
            with gr.Tabs():
                with gr.TabItem('⚙️ Cài đặt'):
                    lecture_pose_style = gr.Slider(minimum=0, maximum=46, step=1, label="Pose style", value=0)
                    lecture_size_of_image = gr.Radio([256, 512], value=256, label='Độ phân giải ảnh')
                    lecture_preprocess_type = gr.Radio(['crop', 'resize','full', 'extcrop', 'extfull'], value='crop')
                    lecture_is_still_mode = gr.Checkbox(label="Still Mode")
                    lecture_batch_size = gr.Slider(label="Batch size", step=1, maximum=10, value=6)
                    lecture_enhancer = gr.Checkbox(label="GFPGAN Face enhancer")
                    lecture_fast_mode_btn = gr.Button('⚡ Chế độ nhanh')
                    lecture_fast_mode_btn.click(
                        fn=set_lecture_fast_mode,
                        outputs=[lecture_size_of_image, lecture_preprocess_type, lecture_is_still_mode, lecture_batch_size, lecture_enhancer, lecture_pose_style]
                    )
            with gr.Tabs():
                gr.Markdown("### 🎬 Video Bài Giảng")
                lecture_final_video = gr.Video(label="Video hoàn chỉnh", format="mp4")
                lecture_info = gr.Textbox(label="Thông tin chi tiết", lines=4, interactive=False)

    extract_lecture_slides_btn.click(
        fn=extract_lecture_slides,
        inputs=[lecture_pptx_file],
        outputs=[lecture_slides_preview]
    )

    # Kết nối generate button với handler
    def handle_generate_lecture(pptx_file, source_image, voice_id, preprocess_type, is_still_mode, enhancer, batch_size, size_of_image, pose_style):
        """Handler cho việc tạo video bài giảng"""
        try:
            if not source_image:
                return "❌ Vui lòng chọn ảnh giáo viên", None, "❌ Thiếu ảnh giáo viên"
            
            if not pptx_file:
                return "❌ Vui lòng chọn file PowerPoint", None, "❌ Thiếu file PowerPoint"
            
            if not voice_id:
                return "❌ Vui lòng chọn giọng nhân bản", None, "❌ Thiếu giọng nhân bản"
            
            # Kiểm tra voice có tồn tại không
            if not has_voice("current_user", voice_id):
                return f"❌ Giọng nói '{voice_id}' không tồn tại. Vui lòng đăng ký lại!", None, f"❌ Voice '{voice_id}' không tồn tại"
            
            # Lấy text từ slides preview (cần implement)
            slides_text = "Hello everybody, I will be teach you..."  # Placeholder
            
            # Tạo audio với voice đã chọn
            audio_path = convert_text_to_audio_with_voice(slides_text, "current_user", voice_id)
            
            if not audio_path:
                return "❌ Không thể tạo audio", None, "❌ Lỗi tạo audio"
            
            # TODO: Implement SadTalker pipeline
            # result_video = generate_sadtalker_video(source_image, audio_path, ...)
            
            return "✅ Đang tạo video...", None, f"✅ Audio đã tạo: {audio_path}"
            
        except Exception as e:
            error_msg = f"❌ Lỗi: {str(e)}"
            return error_msg, None, error_msg

    # Kết nối generate button
    generate_lecture_btn.click(
        fn=handle_generate_lecture,
        inputs=[
            lecture_pptx_file,
            lecture_source_image,
            lecture_voice_id,
            lecture_preprocess_type,
            lecture_is_still_mode,
            lecture_enhancer,
            lecture_batch_size,
            lecture_size_of_image,
            lecture_pose_style
        ],
        outputs=[
            lecture_status,
            lecture_final_video,
            lecture_info
        ]
    )

    return {
        'source_image': lecture_source_image,
        'pptx_file': lecture_pptx_file,
        'voice_id': lecture_voice_id,
        'slides_preview': lecture_slides_preview,
        'extract_btn': extract_lecture_slides_btn,
        'generate_btn': generate_lecture_btn,
        'status': lecture_status,
        'pose_style': lecture_pose_style,
        'size_of_image': lecture_size_of_image,
        'preprocess_type': lecture_preprocess_type,
        'is_still_mode': lecture_is_still_mode,
        'batch_size': lecture_batch_size,
        'enhancer': lecture_enhancer,
        'fast_mode_btn': lecture_fast_mode_btn,
        'final_video': lecture_final_video,
        'info': lecture_info
    }
