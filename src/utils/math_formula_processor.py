import re
import unicodedata
import xml.etree.ElementTree as ET
from typing import List, Dict, Tuple, Optional
import logging

# Thiết lập logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class MathFormulaProcessor:
    """
    Xử lý các ký tự đặc biệt và công thức toán học từ PowerPoint
    """
    
    def __init__(self):
        # Bảng tra cứu ký tự đặc biệt sang tiếng Việt
        self.special_char_map = {
            # Số mũ
            '²': ' mũ hai',
            '³': ' mũ ba',
            '¹': ' mũ một',
            '⁴': ' mũ bốn',
            '⁵': ' mũ năm',
            '⁶': ' mũ sáu',
            '⁷': ' mũ bảy',
            '⁸': ' mũ tám',
            '⁹': ' mũ chín',
            '⁰': ' mũ không',
            
            # Chỉ số dưới
            '₁': ' chỉ số một',
            '₂': ' chỉ số hai',
            '₃': ' chỉ số ba',
            '₄': ' chỉ số bốn',
            '₅': ' chỉ số năm',
            '₆': ' chỉ số sáu',
            '₇': ' chỉ số bảy',
            '₈': ' chỉ số tám',
            '₉': ' chỉ số chín',
            '₀': ' chỉ số không',
            
            # Ký tự Hy Lạp
            'α': ' alpha',
            'β': ' beta',
            'γ': ' gamma',
            'δ': ' delta',
            'ε': ' epsilon',
            'ζ': ' zeta',
            'η': ' eta',
            'θ': ' theta',
            'ι': ' iota',
            'κ': ' kappa',
            'λ': ' lambda',
            'μ': ' mu',
            'ν': ' nu',
            'ξ': ' xi',
            'ο': ' omicron',
            'π': ' pi',
            'ρ': ' rho',
            'σ': ' sigma',
            'τ': ' tau',
            'υ': ' upsilon',
            'φ': ' phi',
            'χ': ' chi',
            'ψ': ' psi',
            'ω': ' omega',
            
            # Chữ hoa Hy Lạp
            'Α': ' Alpha',
            'Β': ' Beta',
            'Γ': ' Gamma',
            'Δ': ' Delta',
            'Ε': ' Epsilon',
            'Ζ': ' Zeta',
            'Η': ' Eta',
            'Θ': ' Theta',
            'Ι': ' Iota',
            'Κ': ' Kappa',
            'Λ': ' Lambda',
            'Μ': ' Mu',
            'Ν': ' Nu',
            'Ξ': ' Xi',
            'Ο': ' Omicron',
            'Π': ' Pi',
            'Ρ': ' Rho',
            'Σ': ' Sigma',
            'Τ': ' Tau',
            'Υ': ' Upsilon',
            'Φ': ' Phi',
            'Χ': ' Chi',
            'Ψ': ' Psi',
            'Ω': ' Omega',
            
            # Ký tự toán học
            '±': ' cộng trừ',
            '∓': ' trừ cộng',
            '×': ' nhân',
            '÷': ' chia',
            '⋅': ' nhân',
            '∗': ' nhân',
            '√': ' căn bậc hai',
            '∛': ' căn bậc ba',
            '∜': ' căn bậc bốn',
            '∞': ' vô cùng',
            '≈': ' xấp xỉ',
            '≠': ' khác',
            '≤': ' nhỏ hơn hoặc bằng',
            '≥': ' lớn hơn hoặc bằng',
            '≪': ' nhỏ hơn rất nhiều',
            '≫': ' lớn hơn rất nhiều',
            '≡': ' đồng dư',
            '≅': ' đồng dạng',
            '∝': ' tỷ lệ thuận',
            '∑': ' tổng',
            '∏': ' tích',
            '∫': ' tích phân',
            '∬': ' tích phân kép',
            '∭': ' tích phân ba',
            '∮': ' tích phân đường',
            '∯': ' tích phân mặt',
            '∰': ' tích phân thể tích',
            '∇': ' nabla',
            '∂': ' đạo hàm riêng',
            '∆': ' delta',
            '∅': ' tập rỗng',
            '∈': ' thuộc',
            '∉': ' không thuộc',
            '∋': ' chứa',
            '∌': ' không chứa',
            '⊂': ' tập con',
            '⊃': ' tập cha',
            '⊆': ' tập con hoặc bằng',
            '⊇': ' tập cha hoặc bằng',
            '∪': ' hợp',
            '∩': ' giao',
            '∖': ' hiệu',
            '⊕': ' tổng trực tiếp',
            '⊗': ' tích tensor',
            '⊥': ' vuông góc',
            '∥': ' song song',
            '∠': ' góc',
            '∡': ' góc đo',
            '∢': ' góc phẳng',
            '°': ' độ',
            '′': ' phút',
            '″': ' giây',
            '‰': ' phần nghìn',
            '‱': ' phần vạn',
            
            # Mũi tên
            '→': ' mũi tên phải',
            '←': ' mũi tên trái',
            '↑': ' mũi tên lên',
            '↓': ' mũi tên xuống',
            '↔': ' mũi tên hai chiều',
            '↕': ' mũi tên lên xuống',
            '⇒': ' suy ra',
            '⇐': ' ngược lại',
            '⇔': ' tương đương',
            '⇎': ' không tương đương',
            
            # Logic
            '∀': ' với mọi',
            '∃': ' tồn tại',
            '∄': ' không tồn tại',
            '∴': ' do đó',
            '∵': ' vì',
            '∧': ' và',
            '∨': ' hoặc',
            '¬': ' không',
            '⊤': ' đúng',
            '⊥': ' sai',
            
            # Tập số
            'ℕ': ' tập số tự nhiên',
            'ℤ': ' tập số nguyên',
            'ℚ': ' tập số hữu tỷ',
            'ℝ': ' tập số thực',
            'ℂ': ' tập số phức',
            'ℙ': ' tập số nguyên tố',
            
            # Ký tự khác
            'ℵ': ' aleph',
            'ℶ': ' beth',
            'ℷ': ' gimel',
            'ℸ': ' daleth',
            'ℏ': ' h bar',
            'ℯ': ' e',
            'ℊ': ' g',
            'ℴ': ' o',
            'ℵ': ' aleph',
        }
        
        # Mẫu regex để nhận diện các công thức toán học đơn giản
        self.math_patterns = [
            # Phân số: a/b
            (r'(\d+)/(\d+)', r'\1 chia \2'),
            # Căn bậc hai: √x (cần xử lý trước khi thay thế √)
            (r'√(\w+)', r'căn bậc hai của \1'),
            # Căn bậc n: n√x (cần xử lý trước khi thay thế √)
            (r'(\d+)√(\w+)', r'\1 căn bậc \1 của \2'),
            # Lũy thừa: x^n
            (r'(\w+)\^(\d+)', r'\1 mũ \2'),
            # Tích phân: ∫f(x)dx
            (r'∫([^d]+)d([a-z])', r'tích phân của \1 theo \2'),
            # Đạo hàm: d/dx
            (r'd/(d[a-z])', r'đạo hàm theo \1'),
            # Tổng: Σ
            (r'Σ([^=]+)=([^=]+)', r'tổng của \1 từ \2'),
            # Tích: Π
            (r'Π([^=]+)=([^=]+)', r'tích của \1 từ \2'),
        ]
    
    def process_special_characters(self, text: str) -> str:
        """
        Chuyển đổi các ký tự đặc biệt thành văn bản tiếng Việt
        """
        if not text:
            return text
        
        processed_text = text
        
        # Bước 1: Xử lý các mẫu toán học phức tạp TRƯỚC
        for pattern, replacement in self.math_patterns:
            processed_text = re.sub(pattern, replacement, processed_text)
        
        # Bước 2: Xử lý từng ký tự đặc biệt
        for special_char, replacement in self.special_char_map.items():
            processed_text = processed_text.replace(special_char, replacement)
        
        # Bước 3: Xử lý các ký tự Unicode khác
        processed_text = self._process_unicode_chars(processed_text)
        
        # Bước 4: Làm sạch văn bản
        processed_text = self._clean_text(processed_text)
        
        return processed_text
    
    def debug_process(self, text: str) -> Dict[str, str]:
        """
        Debug từng bước xử lý để kiểm tra
        """
        debug_info = {
            'original': text,
            'after_regex': text,
            'after_special_chars': text,
            'after_unicode': text,
            'after_cleaning': text,
            'final': text
        }
        
        # Bước 1: Regex
        processed_text = text
        for pattern, replacement in self.math_patterns:
            processed_text = re.sub(pattern, replacement, processed_text)
        debug_info['after_regex'] = processed_text
        
        # Bước 2: Ký tự đặc biệt
        for special_char, replacement in self.special_char_map.items():
            processed_text = processed_text.replace(special_char, replacement)
        debug_info['after_special_chars'] = processed_text
        
        # Bước 3: Unicode
        processed_text = self._process_unicode_chars(processed_text)
        debug_info['after_unicode'] = processed_text
        
        # Bước 4: Làm sạch
        processed_text = self._clean_text(processed_text)
        debug_info['after_cleaning'] = processed_text
        
        debug_info['final'] = processed_text
        return debug_info
    
    def _process_unicode_chars(self, text: str) -> str:
        """
        Xử lý các ký tự Unicode khác
        """
        processed_chars = []
        
        for char in text:
            if ord(char) > 127:  # Ký tự không phải ASCII
                try:
                    # Lấy tên Unicode
                    char_name = unicodedata.name(char)
                    
                    # Xử lý các trường hợp đặc biệt
                    if 'SUPERSCRIPT' in char_name:
                        # Số mũ
                        if char_name.endswith('TWO'):
                            processed_chars.append(' mũ hai')
                        elif char_name.endswith('THREE'):
                            processed_chars.append(' mũ ba')
                        elif char_name.endswith('ONE'):
                            processed_chars.append(' mũ một')
                        else:
                            # Lấy số từ tên
                            number = char_name.split()[-1]
                            processed_chars.append(f' mũ {number}')
                    elif 'SUBSCRIPT' in char_name:
                        # Chỉ số dưới
                        if char_name.endswith('TWO'):
                            processed_chars.append(' chỉ số hai')
                        elif char_name.endswith('THREE'):
                            processed_chars.append(' chỉ số ba')
                        elif char_name.endswith('ONE'):
                            processed_chars.append(' chỉ số một')
                        else:
                            number = char_name.split()[-1]
                            processed_chars.append(f' chỉ số {number}')
                    elif 'GREEK' in char_name:
                        # Chữ Hy Lạp
                        greek_name = char_name.split()[-1].lower()
                        processed_chars.append(f' {greek_name}')
                    elif 'MATHEMATICAL' in char_name:
                        # Ký tự toán học
                        math_name = char_name.split()[-1].lower()
                        processed_chars.append(f' {math_name}')
                    else:
                        # Giữ nguyên ký tự nếu không xử lý được
                        processed_chars.append(char)
                        
                except ValueError:
                    # Nếu không lấy được tên Unicode, giữ nguyên
                    processed_chars.append(char)
            else:
                processed_chars.append(char)
        
        return ''.join(processed_chars)
    
    def _clean_text(self, text: str) -> str:
        """
        Làm sạch văn bản sau khi xử lý
        """
        # Loại bỏ khoảng trắng thừa
        text = re.sub(r'\s+', ' ', text)
        
        # Loại bỏ khoảng trắng đầu cuối
        text = text.strip()
        
        # Xử lý các dấu câu - đảm bảo không có khoảng trắng trước dấu câu
        text = re.sub(r'\s+([.,;:!?])', r'\1', text)
        
        # Xử lý dấu ngoặc - đảm bảo khoảng trắng phù hợp
        text = re.sub(r'\(\s+', '(', text)
        text = re.sub(r'\s+\)', ')', text)
        
        # KHÔNG xử lý khoảng trắng giữa các từ đã được xử lý
        # Vì có thể làm hỏng các từ như "mũ hai", "chỉ số ba"
        
        return text
    
    def extract_math_objects_from_pptx(self, pptx_file_path: str) -> List[Dict]:
        """
        Trích xuất các đối tượng toán học từ PowerPoint
        """
        try:
            from pptx import Presentation
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import qn
            
            prs = Presentation(pptx_file_path)
            math_objects = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_math = []
                
                for shape in slide.shapes:
                    # Kiểm tra các đối tượng toán học
                    if hasattr(shape, 'element'):
                        element = shape.element
                        
                        # Tìm các phần tử toán học
                        math_elements = element.findall('.//m:oMath', {'m': 'http://schemas.openxmlformats.org/officeDocument/2006/math'})
                        
                        for math_elem in math_elements:
                            try:
                                # Chuyển đổi XML thành văn bản
                                math_text = self._extract_mathml_text(math_elem)
                                if math_text:
                                    slide_math.append({
                                        'type': 'math_object',
                                        'content': math_text,
                                        'processed_content': self.process_special_characters(math_text)
                                    })
                            except Exception as e:
                                logger.warning(f"Lỗi xử lý phần tử toán học: {e}")
                
                if slide_math:
                    math_objects.append({
                        'slide_number': slide_num + 1,
                        'math_objects': slide_math
                    })
            
            return math_objects
            
        except ImportError:
            logger.error("Không thể import python-pptx. Vui lòng cài đặt: pip install python-pptx")
            return []
        except Exception as e:
            logger.error(f"Lỗi trích xuất đối tượng toán học: {e}")
            return []
    
    def _extract_mathml_text(self, math_element) -> str:
        """
        Trích xuất văn bản từ phần tử MathML
        """
        try:
            # Chuyển đổi XML thành chuỗi
            xml_str = ET.tostring(math_element, encoding='unicode')
            
            # Xử lý các thẻ MathML cơ bản
            text_parts = []
            
            # Tìm tất cả các phần tử văn bản
            for elem in math_element.iter():
                if elem.text and elem.text.strip():
                    text_parts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    text_parts.append(elem.tail.strip())
            
            return ' '.join(text_parts)
            
        except Exception as e:
            logger.warning(f"Lỗi trích xuất MathML: {e}")
            return ""
    
    def process_powerpoint_text(self, pptx_file_path: str) -> Dict:
        """
        Xử lý toàn bộ văn bản từ PowerPoint, bao gồm cả đối tượng toán học
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_file_path)
            processed_slides = []
            
            # Trích xuất đối tượng toán học
            math_objects = self.extract_math_objects_from_pptx(pptx_file_path)
            math_dict = {obj['slide_number']: obj['math_objects'] for obj in math_objects}
            
            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                text_chunks = []
                
                # Trích xuất văn bản thông thường
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_chunks.append(shape.text.strip())
                
                # Kết hợp văn bản thông thường và toán học
                slide_text = " ".join(filter(None, text_chunks))
                
                # Xử lý văn bản thông thường
                processed_text = self.process_special_characters(slide_text)
                
                # Thêm đối tượng toán học nếu có
                if slide_num in math_dict:
                    math_texts = []
                    for math_obj in math_dict[slide_num]:
                        math_texts.append(math_obj['processed_content'])
                    
                    if math_texts:
                        processed_text += " " + " ".join(math_texts)
                
                processed_slides.append({
                    'slide_number': slide_num,
                    'original_text': slide_text,
                    'processed_text': processed_text,
                    'has_math_objects': slide_num in math_dict
                })
            
            return {
                'slides': processed_slides,
                'total_slides': len(processed_slides),
                'slides_with_math': len([s for s in processed_slides if s['has_math_objects']])
            }
            
        except Exception as e:
            logger.error(f"Lỗi xử lý PowerPoint: {e}")
            return {
                'slides': [],
                'total_slides': 0,
                'slides_with_math': 0,
                'error': str(e)
            }

# Hàm tiện ích để sử dụng nhanh
def process_math_text(text: str) -> str:
    """
    Hàm tiện ích để xử lý nhanh văn bản chứa công thức toán học
    """
    processor = MathFormulaProcessor()
    return processor.process_special_characters(text)

def process_powerpoint_file(pptx_file_path: str) -> Dict:
    """
    Hàm tiện ích để xử lý toàn bộ file PowerPoint
    """
    processor = MathFormulaProcessor()
    return processor.process_powerpoint_text(pptx_file_path)

